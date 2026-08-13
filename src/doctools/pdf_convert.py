"""PDF → Word / PDF → PPT。

- PDF → Word：优先 pdf2docx 版式还原（有损）；引擎失败时回退到
  PyMuPDF 文字提取生成极简 docx（段落 + 表格，无图片/字体还原），
  参照飞鼠的降级链设计；纯扫描件（全页无文字）报 ``PDF_NO_TEXT``
  （OCR 回退见后续规划）。
- PDF → PPT：每页 PDF 渲染成高清图片，插入一张幻灯片（版式还原、文字不可编辑）。

模块顶部不导入三方库，worker 内惰性加载，避免拖慢 CLI 启动。
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from doctools.errors import PDF_CONVERT_ENGINE_FAILED, PDF_NO_TEXT, DoctoolsError
from doctools.model import FileResult, ProgressFn
from doctools.resource_policy import assert_pdf_pages, assert_pixmap_size


def _page_text_rows(page: Any) -> list[str]:
    """返回页面非空文本行（按阅读顺序）。"""
    return [line for line in page.get_text().splitlines() if line.strip()]


def _convert_with_pdf2docx(src: Path, dst: Path) -> None:
    from pdf2docx import Converter  # noqa: PLC0415

    converter = Converter(str(src))
    try:
        converter.convert(str(dst))
    finally:
        converter.close()


def _write_fallback_docx(src: Path, dst: Path) -> None:
    """用 PyMuPDF 文本 + 表格检测生成极简 docx（回退产物）。

    每页：页标题 + 表格区域外的文本块逐行成段；检测到的表格按
    extract() 行写入 docx 表格。参照飞鼠 fallback 的"多列行→表格、
    单列行→段落"思路，但表格结构来自 PyMuPDF 的 find_tables。
    """
    import fitz  # noqa: PLC0415  # PyMuPDF
    from docx import Document  # noqa: PLC0415

    document = Document()
    doc = fitz.open(str(src))
    try:
        for index, page in enumerate(doc, start=1):
            document.add_heading(f"Page {index}", level=1)
            tables = page.find_tables()
            table_rects = [tuple(table.bbox) for table in tables.tables] if tables else []
            for block in page.get_text("blocks"):
                x0, y0, x1, y1 = block[0], block[1], block[2], block[3]
                text = str(block[4]).strip()
                if not text:
                    continue
                # 完全落在某个表格 bbox 内的文本块由表格呈现，跳过
                if any(
                    x0 >= rx0 - 2 and y0 >= ry0 - 2 and x1 <= rx1 + 2 and y1 <= ry1 + 2
                    for (rx0, ry0, rx1, ry1) in table_rects
                ):
                    continue
                for line in text.splitlines():
                    if line.strip():
                        document.add_paragraph(line.strip())
            for table in tables.tables:
                rows = table.extract()
                if not rows:
                    continue
                ncols = max(len(r) for r in rows)
                docx_table = document.add_table(rows=0, cols=ncols)
                for row in rows:
                    cells = docx_table.add_row().cells
                    for i in range(ncols):
                        cells[i].text = row[i] if i < len(row) and row[i] is not None else ""
    finally:
        doc.close()
    document.save(str(dst))


def pdf_to_docx(src: Path, dst: Path) -> str | None:
    """把单个 PDF 转成 Word（作为 process_batch 的 worker 使用）。

    返回附注（如回退说明），无附注时返回 None。全页无文字（扫描件）
    抛 ``PDF_NO_TEXT``。
    """
    dst.parent.mkdir(parents=True, exist_ok=True)
    import fitz  # noqa: PLC0415  # PyMuPDF

    doc = fitz.open(str(src))
    try:
        assert_pdf_pages(len(doc))
        has_text = any(_page_text_rows(page) for page in doc)
    finally:
        doc.close()
    if not has_text:
        raise DoctoolsError(
            PDF_NO_TEXT,
            "这个 PDF 没有可提取的文字，可能是扫描版图片 PDF。"
            "请先对扫描件做 OCR 后再转换（OCR 支持规划中）。",
            "This PDF has no extractable text; it may be a scanned image PDF. "
            "OCR support is planned.",
        )

    try:
        _convert_with_pdf2docx(src, dst)
    except Exception as exc:  # noqa: BLE001 - pdf2docx 失败即回退文字提取
        try:
            _write_fallback_docx(src, dst)
        except Exception as fallback_exc:  # noqa: BLE001 - 回退也失败才报错
            raise DoctoolsError(
                PDF_CONVERT_ENGINE_FAILED,
                f"PDF 转 Word 失败：pdf2docx 与文字提取回退均未成功。"
                f"\npdf2docx：{exc}\n文字提取：{fallback_exc}",
                f"PDF to Word failed: both pdf2docx and text-extraction fallback failed. "
                f"\npdf2docx: {exc}\ntext extraction: {fallback_exc}",
            ) from fallback_exc
        return "pdf2docx 版式还原失败，已回退为文字提取（无图片/字体还原）。"
    return None


def pdf_to_pptx(src: Path, dst: Path) -> None:
    """把单个 PDF 转成 PPT：每页渲染成图片，插入一张幻灯片（铺满整页）。"""
    dst.parent.mkdir(parents=True, exist_ok=True)
    import io  # noqa: PLC0415

    import fitz  # noqa: PLC0415  # PyMuPDF
    from pptx import Presentation  # noqa: PLC0415

    prs = Presentation()
    doc = fitz.open(str(src))
    try:
        assert_pdf_pages(len(doc))
        for index, page in enumerate(doc):
            if index == 0:
                # 幻灯片尺寸取第一页大小（PDF 单位是 pt，1pt = 12700 EMU）
                prs.slide_width = int(page.rect.width * 12700)
                prs.slide_height = int(page.rect.height * 12700)
            pix = page.get_pixmap(dpi=150)
            assert_pixmap_size(pix.width, pix.height)
            image = io.BytesIO(pix.tobytes("png"))
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
            slide.shapes.add_picture(image, 0, 0, width=prs.slide_width, height=prs.slide_height)
    finally:
        doc.close()
    prs.save(str(dst))


def pdf_to_images(
    src: Path,
    out_dir: Path,
    on_progress: ProgressFn | None = None,
) -> list[FileResult]:
    """把 PDF 的每一页渲染成一张 PNG 图片，命名 ``{stem}_p{n}.png``。"""
    import fitz  # noqa: PLC0415  # PyMuPDF

    out_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(src))
    results: list[FileResult] = []
    try:
        assert_pdf_pages(len(doc))
        total = len(doc)
        for index, page in enumerate(doc, start=1):
            pix = page.get_pixmap(dpi=150)
            assert_pixmap_size(pix.width, pix.height)
            dst = out_dir / f"{src.stem}_p{index}.png"
            pix.save(str(dst))
            result = FileResult(src=src, dst=dst, ok=True)
            results.append(result)
            if on_progress is not None:
                on_progress(total, index, result)
    finally:
        doc.close()
    return results
