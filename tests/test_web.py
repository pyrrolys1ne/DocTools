"""Web 后端 API 冒烟测试。

未安装 fastapi 时整模块跳过（``.[dev]`` 裸装也能跑核心测试）。
"""

from __future__ import annotations

import time
from pathlib import Path

import pytest
from pypdf import PdfReader, PdfWriter

pytest.importorskip("fastapi")

from fastapi.testclient import TestClient  # noqa: E402
from test_docx import _doc_with_headers  # noqa: E402

from web.app import app  # noqa: E402

client = TestClient(app)


def _wait_done(job_id: str, timeout: float = 10.0) -> dict:
    deadline = time.time() + timeout
    while time.time() < deadline:
        data = client.get(f"/api/jobs/{job_id}").json()
        if data["status"] in ("done", "failed"):
            return data
        time.sleep(0.05)
    raise TimeoutError(f"任务 {job_id} 超时未完成")


def test_scan_lists_docx_recursively(tmp_path: Path) -> None:
    (tmp_path / "sub").mkdir()
    (tmp_path / "a.txt").write_text("hi")
    (tmp_path / "a.docx").write_bytes(b"x")
    (tmp_path / "sub" / "b.docx").write_bytes(b"x")

    resp = client.post(
        "/api/scan", json={"source_path": str(tmp_path), "recursive": True}
    )

    assert resp.status_code == 200
    assert resp.json()["kind"] == "dir"
    names = [f["name"] for f in resp.json()["files"]]
    assert names == ["a.docx", "sub/b.docx"]


def test_drives_lists_at_least_one() -> None:
    resp = client.get("/api/drives")
    assert resp.status_code == 200
    data = resp.json()
    assert isinstance(data["drives"], list)
    assert len(data["drives"]) >= 1
    assert isinstance(data["special"], list)


def test_drives_include_special_folders(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """常见用户目录（桌面/文档…）应解析为可访问的物理路径。"""
    (tmp_path / "Desktop").mkdir()
    (tmp_path / "Documents").mkdir()
    monkeypatch.setattr("web.app.Path.home", staticmethod(lambda: tmp_path))

    resp = client.get("/api/drives")

    assert resp.status_code == 200
    special = resp.json()["special"]
    assert {"name": "桌面", "path": str(tmp_path / "Desktop")} in special
    assert {"name": "文档", "path": str(tmp_path / "Documents")} in special
    assert "视频" not in {s["name"] for s in special}  # 未创建的目录不列出


def test_explore_returns_parent_and_dirs(tmp_path: Path) -> None:
    (tmp_path / "sub").mkdir()
    (tmp_path / "a.txt").write_text("hi")
    (tmp_path / "a.docx").write_bytes(b"x")

    resp = client.get(f"/api/explore?dir={tmp_path}")

    assert resp.status_code == 200
    data = resp.json()
    assert data["dirs"] == ["sub"]
    assert data["files"] == ["a.docx"]
    assert data["parent"] == str(tmp_path.resolve().parent)
    assert data["error"] is None


def test_explore_hides_dollar_prefixed_dirs(tmp_path: Path) -> None:
    """以 $ 开头的系统特殊文件夹（如 $RECYCLE.BIN）不应出现在可选项里。"""
    (tmp_path / "$RECYCLE.BIN").mkdir()
    (tmp_path / "sub").mkdir()

    resp = client.get(f"/api/explore?dir={tmp_path}")

    assert resp.status_code == 200
    assert resp.json()["dirs"] == ["sub"]


def test_explore_handles_unreadable_dir(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    """权限受限的系统目录（如 $RECYCLE.BIN 子层）不应 500，而是返回提示。"""
    import pathlib

    locked = tmp_path / "locked"
    locked.mkdir()
    real_iterdir = pathlib.Path.iterdir

    def fake_iterdir(self: Path):
        if self == locked:
            raise PermissionError("denied")
        return real_iterdir(self)

    monkeypatch.setattr(pathlib.Path, "iterdir", fake_iterdir)

    resp = client.get(f"/api/explore?dir={locked}")

    assert resp.status_code == 200
    data = resp.json()
    assert data["dirs"] == []
    assert data["files"] == []
    assert data["error"] is not None


def test_explore_404_on_missing_dir(tmp_path: Path) -> None:
    resp = client.get(f"/api/explore?dir={tmp_path / 'nope'}")
    assert resp.status_code == 404


def test_scan_404_on_missing_dir(tmp_path: Path) -> None:
    resp = client.post("/api/scan", json={"source_path": str(tmp_path / "nope")})
    assert resp.status_code == 404


def test_explore_filters_files_by_ext(tmp_path: Path) -> None:
    (tmp_path / "a.docx").write_bytes(b"x")
    (tmp_path / "b.pdf").write_bytes(b"x")
    (tmp_path / "c.pptx").write_bytes(b"x")

    resp = client.get(f"/api/explore?dir={tmp_path}&exts=.pdf")

    assert resp.status_code == 200
    assert resp.json()["files"] == ["b.pdf"]

    resp2 = client.get(f"/api/explore?dir={tmp_path}&exts=.docx,.pptx")
    assert resp2.json()["files"] == ["a.docx", "c.pptx"]


def test_scan_single_file(tmp_path: Path) -> None:
    _doc_with_headers().save(str(tmp_path / "a.docx"))

    resp = client.post("/api/scan", json={"source_path": str(tmp_path / "a.docx")})

    assert resp.status_code == 200
    data = resp.json()
    assert data["kind"] == "file"
    assert data["files"] == [{"name": "a.docx", "size": (tmp_path / "a.docx").stat().st_size}]


def test_job_dry_run(tmp_path: Path) -> None:
    src = tmp_path / "in"
    src.mkdir()
    _doc_with_headers().save(str(src / "a.docx"))

    resp = client.post(
        "/api/jobs",
        json={
            "source_path": str(src),
            "output_path": str(tmp_path / "out"),
            "dry_run": True,
        },
    )
    job_id = resp.json()["id"]

    data = _wait_done(job_id)
    assert data["status"] == "done"
    assert data["total"] == 1
    assert data["done"] == 1
    assert not (tmp_path / "out" / "a.docx").exists()  # dry-run 不写文件


def test_job_processing(tmp_path: Path) -> None:
    src = tmp_path / "in"
    src.mkdir()
    _doc_with_headers().save(str(src / "a.docx"))

    resp = client.post(
        "/api/jobs",
        json={"source_path": str(src), "output_path": str(tmp_path / "out")},
    )
    job_id = resp.json()["id"]

    data = _wait_done(job_id)
    assert data["status"] == "done"
    assert data["results"][0]["ok"] is True
    assert (tmp_path / "out" / "a.docx").exists()


def test_job_single_file(tmp_path: Path) -> None:
    _doc_with_headers().save(str(tmp_path / "a.docx"))

    resp = client.post(
        "/api/jobs",
        json={"source_path": str(tmp_path / "a.docx")},
    )
    job_id = resp.json()["id"]

    data = _wait_done(job_id)
    assert data["status"] == "done"
    assert data["total"] == 1
    assert data["results"][0]["ok"] is True
    # 未指定输出路径时，默认生成 {stem}_cleaned.docx
    assert (tmp_path / "a_cleaned.docx").exists()


def test_job_single_file_output_is_dir(tmp_path: Path) -> None:
    _doc_with_headers().save(str(tmp_path / "a.docx"))
    out_dir = tmp_path / "out"

    resp = client.post(
        "/api/jobs",
        json={
            "source_path": str(tmp_path / "a.docx"),
            "output_path": str(out_dir),
            "output_is_dir": True,
        },
    )
    job_id = resp.json()["id"]

    data = _wait_done(job_id)
    assert data["status"] == "done"
    assert data["results"][0]["ok"] is True
    assert (out_dir / "a_cleaned.docx").exists()


def test_job_ws_streams_progress(tmp_path: Path) -> None:
    src = tmp_path / "in"
    src.mkdir()
    _doc_with_headers().save(str(src / "a.docx"))

    resp = client.post(
        "/api/jobs",
        json={"source_path": str(src), "output_path": str(tmp_path / "out")},
    )
    job_id = resp.json()["id"]

    with client.websocket_connect(f"/api/jobs/{job_id}/ws") as ws:
        data = ws.receive_json()
        assert data["id"] == job_id
        assert data["status"] in ("pending", "running", "done", "failed")


# ---------- PDF 合并 / 拆分 / 转 PDF ----------


def _make_pdf(path: Path, pages: int) -> None:
    writer = PdfWriter()
    for _ in range(pages):
        writer.add_blank_page(width=200, height=200)
    with path.open("wb") as f:
        writer.write(f)


def test_job_merge_pdf(tmp_path: Path) -> None:
    a, b = tmp_path / "a.pdf", tmp_path / "b.pdf"
    _make_pdf(a, 2)
    _make_pdf(b, 3)
    merged = tmp_path / "merged.pdf"

    resp = client.post(
        "/api/jobs",
        json={
            "operation": "merge-pdf",
            "output_path": str(merged),
            "sources": [str(a), str(b)],
        },
    )
    job_id = resp.json()["id"]

    data = _wait_done(job_id)
    assert data["status"] == "done"
    assert data["total"] == 2
    assert merged.exists()
    assert len(PdfReader(str(merged)).pages) == 5


def test_job_merge_requires_sources(tmp_path: Path) -> None:
    resp = client.post(
        "/api/jobs",
        json={"operation": "merge-pdf", "output_path": str(tmp_path / "m.pdf")},
    )
    job_id = resp.json()["id"]
    data = _wait_done(job_id)
    assert data["status"] == "failed"
    assert "源文件" in data["error"]


def test_job_split_pdf_ranges(tmp_path: Path) -> None:
    src = tmp_path / "in.pdf"
    _make_pdf(src, 5)
    out = tmp_path / "out"

    resp = client.post(
        "/api/jobs",
        json={
            "operation": "split-pdf",
            "source_path": str(src),
            "output_path": str(out),
            "page_ranges": "1-2,4",
        },
    )
    job_id = resp.json()["id"]

    data = _wait_done(job_id)
    assert data["status"] == "done"
    assert sorted(Path(r["dst"]).name for r in data["results"]) == ["in_p1-2.pdf", "in_p4.pdf"]
    assert (out / "in_p1-2.pdf").exists()


def test_job_split_pdf_every_page(tmp_path: Path) -> None:
    src = tmp_path / "in.pdf"
    _make_pdf(src, 3)
    out = tmp_path / "out"

    resp = client.post(
        "/api/jobs",
        json={"operation": "split-pdf", "source_path": str(src), "output_path": str(out)},
    )
    job_id = resp.json()["id"]

    data = _wait_done(job_id)
    assert data["status"] == "done"
    assert data["total"] == 3
    assert len(list(out.glob("*.pdf"))) == 3


def test_job_split_requires_file(tmp_path: Path) -> None:
    resp = client.post(
        "/api/jobs",
        json={"operation": "split-pdf", "source_path": str(tmp_path)},
    )
    job_id = resp.json()["id"]
    data = _wait_done(job_id)
    assert data["status"] == "failed"
    assert "单个 PDF 文件" in data["error"]


def _word_available() -> bool:
    """检查本机是否装有 Microsoft Word（不启动 COM，避免测试互相干扰）。"""
    import os  # noqa: PLC0415

    program_files = [
        Path(os.environ.get("ProgramFiles", "C:/Program Files")),
        Path(os.environ.get("ProgramFiles(x86)", "C:/Program Files (x86)")),
    ]
    return any(
        (base / "Microsoft Office" / "root" / "Office16" / "WINWORD.EXE").exists()
        for base in program_files
    )


def test_job_word_to_pdf(tmp_path: Path) -> None:
    if not _word_available():
        pytest.skip("本机未安装 Microsoft Office")
    from docx import Document  # noqa: PLC0415

    doc = Document()
    doc.add_paragraph("hello")
    src = tmp_path / "a.docx"
    doc.save(str(src))
    out = tmp_path / "out"

    resp = client.post(
        "/api/jobs",
        json={
            "operation": "word-to-pdf",
            "source_path": str(src),
            "output_path": str(out),
            "output_is_dir": True,
        },
    )
    job_id = resp.json()["id"]

    data = _wait_done(job_id, timeout=60)
    assert data["status"] == "done"
    assert (out / "a.pdf").exists()


def test_job_ppt_to_pdf(tmp_path: Path) -> None:
    if not _word_available():
        pytest.skip("本机未安装 Microsoft Office")
    from pptx import Presentation  # noqa: PLC0415
    from pptx.util import Inches  # noqa: PLC0415

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2)).text_frame.text = "测试"
    src = tmp_path / "a.pptx"
    prs.save(str(src))
    out = tmp_path / "out"

    resp = client.post(
        "/api/jobs",
        json={
            "operation": "ppt-to-pdf",
            "source_path": str(src),
            "output_path": str(out),
            "output_is_dir": True,
        },
    )
    job_id = resp.json()["id"]

    data = _wait_done(job_id, timeout=60)
    assert data["status"] == "done"
    assert (out / "a.pdf").exists()


def _make_img(path: Path) -> None:
    from PIL import Image  # noqa: PLC0415

    Image.new("RGB", (80, 80), (10, 20, 30)).save(path)


def test_job_image_to_pdf_merges_all(tmp_path: Path) -> None:
    """图片转 PDF 只保留多合一：目录内所有图片合成一个 PDF。"""
    _make_img(tmp_path / "a.png")
    _make_img(tmp_path / "b.png")
    out = tmp_path / "out"

    resp = client.post(
        "/api/jobs",
        json={"operation": "image-to-pdf", "source_path": str(tmp_path), "output_path": str(out)},
    )
    job_id = resp.json()["id"]

    data = _wait_done(job_id)
    assert data["status"] == "done"
    assert data["total"] == 2
    assert (out / "merged.pdf").exists()
    assert len(PdfReader(str(out / "merged.pdf")).pages) == 2
