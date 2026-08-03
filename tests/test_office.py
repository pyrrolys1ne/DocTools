"""docx/pptx → PDF（基于本机 Microsoft Office COM）测试。

未安装 pywin32 或 Microsoft Office 时整模块跳过。
"""

from __future__ import annotations

from pathlib import Path

import pytest

pytest.importorskip("win32com")

from docx import Document  # noqa: E402

from doctools.office import OfficeConverter  # noqa: E402


def _office_available() -> bool:
    """检查本机是否装有 Microsoft Word（不启动 COM，避免测试互相干扰）。"""
    import os  # noqa: PLC0415
    from pathlib import Path  # noqa: PLC0415

    program_files = [
        Path(os.environ.get("ProgramFiles", "C:/Program Files")),
        Path(os.environ.get("ProgramFiles(x86)", "C:/Program Files (x86)")),
    ]
    return any(
        (base / "Microsoft Office" / "root" / "Office16" / "WINWORD.EXE").exists()
        for base in program_files
    )


@pytest.mark.skipif(not _office_available(), reason="本机未安装 Microsoft Office")
def test_docx_to_pdf(tmp_path: Path) -> None:
    doc = Document()
    doc.add_heading("标题", 0)
    doc.add_paragraph("测试内容 hello 123")
    src = tmp_path / "a.docx"
    doc.save(str(src))
    dst = tmp_path / "a.pdf"

    with OfficeConverter() as converter:
        converter.convert(src, dst)

    assert dst.exists()
    assert dst.read_bytes().startswith(b"%PDF")


@pytest.mark.skipif(not _office_available(), reason="本机未安装 Microsoft Office")
def test_pptx_to_pdf(tmp_path: Path) -> None:
    from pptx import Presentation  # noqa: PLC0415
    from pptx.util import Inches  # noqa: PLC0415

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2)).text_frame.text = "测试"
    src = tmp_path / "a.pptx"
    prs.save(str(src))
    dst = tmp_path / "a.pdf"

    with OfficeConverter() as converter:
        converter.convert(src, dst)

    assert dst.exists()
    assert dst.read_bytes().startswith(b"%PDF")
